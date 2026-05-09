"""IMRAD sunumunu programatik olarak üretir (16 slayt, 16:9, yeşil/beyaz tema).

Çıktı: presentation/IMRAD_slides.pptx

Çalıştırma:
    python presentation/build_slides.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "results" / "figures"
OUT = ROOT / "presentation" / "IMRAD_slides.pptx"

# ---------- tema ----------
GREEN_DARK   = RGBColor(0x14, 0x53, 0x2D)
GREEN        = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LIGHT  = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_PALE   = RGBColor(0xF0, 0xFD, 0xF4)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_900    = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700    = RGBColor(0x33, 0x41, 0x55)
SLATE_500    = RGBColor(0x64, 0x74, 0x8B)
SLATE_300    = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200    = RGBColor(0xE2, 0xE8, 0xF0)
AMBER_600    = RGBColor(0xD9, 0x77, 0x06)
BLUE_600     = RGBColor(0x25, 0x63, 0xEB)
RED_600      = RGBColor(0xDC, 0x26, 0x26)

FONT = "Calibri"
TOTAL_SLIDES = 16

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================================
# Yardımcı çizim fonksiyonları
# ============================================================================

def add_textbox(slide, left, top, width, height, text, *,
                font_size=14, bold=False, color=SLATE_700,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name=FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, paragraphs, *,
                   font_size=12, color=SLATE_700, line_spacing=1.35):
    """Çoklu paragraf metni (tam cümleler için)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # bold marker desteği — **kelime** -> kalın
        runs = _split_bold(text)
        for run_text, is_bold in runs:
            r = p.add_run()
            r.text = run_text
            r.font.name = FONT
            r.font.size = Pt(font_size)
            r.font.bold = is_bold
            r.font.color.rgb = SLATE_900 if is_bold else color
    return tb


def _split_bold(text: str) -> list[tuple[str, bool]]:
    """Markdown benzeri **kalın** parser'ı."""
    parts: list[tuple[str, bool]] = []
    i = 0
    while i < len(text):
        if text[i:i+2] == "**":
            end = text.find("**", i + 2)
            if end == -1:
                parts.append((text[i:], False))
                break
            parts.append((text[i+2:end], True))
            i = end + 2
        else:
            j = text.find("**", i)
            if j == -1:
                parts.append((text[i:], False))
                break
            parts.append((text[i:j], False))
            i = j
    return parts


def add_rect(slide, left, top, width, height, fill, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_rounded_rect(slide, left, top, width, height, fill, line=None, line_w=1.0, adj=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = adj
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_oval(slide, left, top, width, height, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.0)
    s.shadow.inherit = False
    return s


def add_line(slide, x1, y1, x2, y2, color=SLATE_300, weight=0.5):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def apply_chrome(slide, page_num: int):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), GREEN)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), GREEN_PALE)
    add_textbox(
        slide, Inches(0.5), SLIDE_H - Inches(0.34), Inches(8), Inches(0.3),
        "YZM304 Derin Öğrenme · Proje 4 · Mehmet Ali Topkara",
        font_size=10, color=GREEN_DARK, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.34), Inches(0.8), Inches(0.3),
        f"{page_num} / {TOTAL_SLIDES}",
        font_size=10, color=GREEN_DARK, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )


def add_slide_title(slide, title: str, subtitle: str | None = None):
    add_textbox(
        slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.7),
        title, font_size=26, bold=True, color=GREEN_DARK,
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.5), Inches(1.02), Inches(12.3), Inches(0.4),
            subtitle, font_size=13, color=SLATE_500,
        )


def add_image(slide, path: Path, left, top, width=None, height=None):
    if not path.exists():
        return add_rect(slide, left, top, width or Inches(4), height or Inches(3),
                        GREEN_PALE, GREEN)
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, height=height)


# ----- diyagram çizimleri ---------------------------------------------------

def _i(v) -> Emu:
    """Float → Emu int (PowerPoint shape coordinatları int olmalı)."""
    return Emu(int(v))


def draw_mlp_diagram(slide, left, top, width, height):
    """MLP'yi 4 katmanlı, daireli ağ diyagramı olarak çiz."""
    layers = [
        ("Girdi", 5, "150,528\n(3×224×224)"),
        ("FC 256", 6, "256 nöron\nReLU + Dropout 0.3"),
        ("FC 128", 5, "128 nöron\nReLU + Dropout 0.3"),
        ("Çıktı", 4, "10 sınıf\n(softmax)"),
    ]
    n = len(layers)
    col_w = _i(width / n)
    neuron_d = Inches(0.36)
    label_h = Inches(0.30)
    sublabel_h = Inches(0.55)

    available_h = _i(height - label_h - sublabel_h - Inches(0.2))
    layer_centers: list[list[tuple[int, int]]] = []

    for i, (lbl, count, sub) in enumerate(layers):
        col_x = _i(left + col_w * i)
        col_center_x = _i(col_x + col_w / 2)

        add_textbox(
            slide, col_x, top, col_w, label_h,
            lbl, font_size=12, bold=True, color=GREEN_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        spacing = _i(available_h / max(count, 1))
        circles_y_start = _i(top + label_h + Inches(0.05))
        centers: list[tuple[int, int]] = []
        for k in range(count):
            cy = _i(circles_y_start + spacing * k + (spacing - neuron_d) / 2)
            cx = _i(col_center_x - neuron_d / 2)
            color = GREEN if i == n - 1 or i == 0 else GREEN_LIGHT
            add_oval(slide, cx, cy, neuron_d, neuron_d, color, GREEN_DARK)
            centers.append((int(col_center_x), int(cy + neuron_d / 2)))
        layer_centers.append(centers)

        add_textbox(
            slide, col_x, _i(top + height - sublabel_h), col_w, sublabel_h,
            sub, font_size=10, color=SLATE_700,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

    for i in range(1, n):
        prev = layer_centers[i - 1]
        cur = layer_centers[i]
        for (x1, y1) in prev:
            for (x2, y2) in cur:
                add_line(
                    slide,
                    _i(x1 + neuron_d / 2), _i(y1),
                    _i(x2 - neuron_d / 2), _i(y2),
                    color=SLATE_300, weight=0.4,
                )


def draw_cnn_diagram(slide, left, top, width, height):
    """CNN scratch'i kademeli azalan dikdörtgen feature map'lerle çiz."""
    stages = [
        ("Girdi",     "3 × 224 × 224", 0.85, GREEN_PALE),
        ("Conv 32",   "32 × 112 × 112", 0.65, GREEN_LIGHT),
        ("Conv 64",   "64 × 56 × 56",   0.50, GREEN_LIGHT),
        ("Conv 128",  "128 × 28 × 28",  0.40, GREEN),
        ("Conv 256",  "256 × 14 × 14",  0.32, GREEN),
        ("GAP+FC",    "256 → 128 → 10", 0.18, GREEN_DARK),
    ]
    n = len(stages)
    col_w = _i(width / n)
    block_w = _i(col_w * 0.55)
    label_h = Inches(0.28)
    sublabel_h = Inches(0.55)

    available_h = _i(height - label_h - sublabel_h - Inches(0.2))

    for i, (lbl, sub, scale, fill) in enumerate(stages):
        col_x = _i(left + col_w * i)
        col_center_x = _i(col_x + col_w / 2)

        add_textbox(
            slide, col_x, top, col_w, label_h,
            lbl, font_size=11, bold=True, color=GREEN_DARK,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        block_h = _i(available_h * scale)
        bx = _i(col_center_x - block_w / 2)
        by = _i(top + label_h + (available_h - block_h) / 2)
        line_color = WHITE if fill == GREEN_DARK else GREEN_DARK
        add_rect(slide, bx, by, block_w, block_h, fill, line_color, line_w=1.2)

        add_textbox(
            slide, col_x, _i(top + height - sublabel_h), col_w, sublabel_h,
            sub, font_size=9, color=SLATE_700,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

        if i < n - 1:
            arrow_y = _i(top + label_h + available_h / 2)
            arrow_x_start = _i(bx + block_w + Inches(0.05))
            arrow_x_end = _i(left + col_w * (i + 1) + (col_w - block_w) / 2 - Inches(0.05))
            arrow_w_int = max(int(arrow_x_end - arrow_x_start), Inches(0.05))
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x_start, _i(arrow_y - Inches(0.07)),
                Emu(arrow_w_int), Inches(0.14),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREEN
            arrow.line.fill.background()


# ============================================================================
# 16 slayt
# ============================================================================

def slide_01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, Inches(2.6), GREEN_DARK)
    add_rect(s, 0, Inches(2.6), SLIDE_W, Inches(0.18), GREEN)

    add_textbox(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.4),
                "YZM304 — DERİN ÖĞRENME · PROJE 4",
                font_size=14, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.9),
                "Oyun Ekran Görüntülerinden Oyun Tespiti",
                font_size=34, bold=True, color=WHITE)
    add_textbox(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
                "MLP, Sıfırdan CNN ve Transfer Learning Mimarilerinin Karşılaştırmalı Analizi",
                font_size=15, color=GREEN_LIGHT)

    add_textbox(s, Inches(0.7), Inches(3.5), Inches(12), Inches(0.5),
                "Mehmet Ali Topkara — 23291093",
                font_size=22, bold=True, color=SLATE_900)
    add_textbox(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
                "Ankara Üniversitesi · Yapay Zeka ve Veri Mühendisliği",
                font_size=15, color=SLATE_700)
    add_textbox(s, Inches(0.7), Inches(4.45), Inches(12), Inches(0.5),
                "2025–2026 Bahar Dönemi",
                font_size=13, color=SLATE_500)

    add_rounded_rect(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.4), GREEN_PALE, GREEN)
    add_textbox(s, Inches(1.0), Inches(5.55), Inches(11.4), Inches(0.4),
                "Proje Özeti",
                font_size=11, bold=True, color=GREEN_DARK)
    add_textbox(s, Inches(1.0), Inches(5.9), Inches(11.4), Inches(0.85),
                "10 sınıflı oyun ekran görüntüsü sınıflandırması · "
                "5 model: 2 baseline (MLP, CNN scratch) + 3 transfer learning "
                "(ResNet50, EfficientNetB0, ViT-Base/16) · "
                "FastAPI + React lokal web demosu · belirsizlik göstergesi",
                font_size=12, color=SLATE_700)

    apply_chrome(s, 1)


def slide_02_motivation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Neden Bu Konu?",
                    "Gaming sektörü, manuel etiketlemenin sınırları ve sınıf çeşitliliği")

    add_paragraphs(
        s, Inches(0.7), Inches(1.7), Inches(12), Inches(2.2),
        [
            "Oyun sektörünün küresel hacmi 2024 itibarıyla yaklaşık **200 milyar dolara** "
            "ulaşmıştır. Twitch ve YouTube Gaming gibi platformlara her gün milyonlarca "
            "saatlik gameplay görüntüsü yüklenmekte, bu içeriğin manuel olarak kategorize "
            "edilmesi pratikte mümkün olmamaktadır.",
            "Otomatik sınıflandırma sistemleri bu noktada arama, içerik moderasyonu, "
            "telif kontrolü ve öneri sistemleri için kritik bir altyapı haline gelmiştir. "
            "Görsel olarak da zorlu bir problemdir: özellikle FPS oyunları benzer "
            "arayüzlere sahiptir, MMORPG'ler aynı sahnede çok farklı görseller üretir, "
            "ve aynı oyun farklı modlarında oldukça farklı görünebilir.",
        ],
        font_size=14,
    )

    # Çeşitlilik notu — vurgulu kutu
    add_rounded_rect(s, Inches(0.7), Inches(4.5), Inches(12), Inches(2.0), GREEN_PALE, GREEN)
    add_textbox(s, Inches(1.0), Inches(4.65), Inches(11.4), Inches(0.4),
                "Konu Seçiminde Çeşitlilik",
                font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(1.0), Inches(5.05), Inches(11.4), Inches(1.4),
        [
            "Sınıf arkadaşlarımın büyük çoğunluğu **sağlık ve insan tabanlı** "
            "görüntü sınıflandırma problemlerini (tıbbi görüntüler, yüz tanıma, "
            "duygu analizi vb.) tercih etmiştir. Ben **çeşitlilik** açısından "
            "farklı bir bağlam seçerek **oyun ekran görüntülerini** seçtim — "
            "hem sektörel olarak güncel, hem de görsel zorluk açısından zengin "
            "bir problem alanı sunduğu için.",
        ],
        font_size=13,
    )

    apply_chrome(s, 2)


def slide_03_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Problem, Araştırma Sorusu ve Birincil Metrik",
                    "Closed-set sınıflandırma · hipotez · neden Macro-F1?")

    # Problem
    add_textbox(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.35),
                "Problem Tanımı", font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.7),
        [
            "Tek bir oyun ekran görüntüsü verildiğinde, modelin görseli **10 önceden "
            "tanımlı sınıftan birine** atfetmesi (closed-set sınıflandırma).",
        ],
        font_size=12,
    )

    # Araştırma sorusu
    add_textbox(s, Inches(0.7), Inches(2.85), Inches(12), Inches(0.35),
                "Araştırma Sorusu", font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.95),
        [
            "Görüntü yapısını yok sayan **MLP**, sıfırdan eğitilen **klasik CNN** "
            "ve ImageNet üzerinde eğitilmiş **transfer learning** modelleri aynı "
            "veri setinde nasıl performans gösterir? Mimari değişiminin ve "
            "pretrained ağırlıkların etkisi sayısal olarak ne kadardır?",
        ],
        font_size=12,
    )

    # Hipotez
    add_textbox(s, Inches(0.7), Inches(4.25), Inches(12), Inches(0.35),
                "Hipotez", font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.7),
        [
            "MLP düşük performans verir; CNN scratch belirgin sıçrama yapar; "
            "transfer learning marjinal ek iyileştirme sağlar — fakat **maliyet "
            "açısından** transfer modeller çok daha pahalıdır.",
        ],
        font_size=12,
    )

    # Birincil metrik kutusu
    add_rounded_rect(s, Inches(0.7), Inches(5.45), Inches(12), Inches(1.5), GREEN_PALE, GREEN)
    add_textbox(s, Inches(1.0), Inches(5.6), Inches(11.4), Inches(0.4),
                "Birincil Metrik: Macro-F1",
                font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(1.0), Inches(6.0), Inches(11.4), Inches(0.95),
        [
            "Veri dengeli (sınıf başına 1000 örnek) olduğundan accuracy ve "
            "weighted-F1 yakın çıkmaktadır. Ancak **Macro-F1** her sınıfa "
            "**eşit ağırlık** verir; precision ve recall'u harmonik ortalama "
            "ile birleştirir — bu da adil mimari karşılaştırması için "
            "akademik standarttır. Tüm yorumlar bu metrik üzerinden yapılır.",
        ],
        font_size=11,
    )

    apply_chrome(s, 3)


def slide_04_dataset(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Veri Seti — Kaynak ve Yapısı",
                    "Kaggle Gameplay Images · 10 sınıf × 1000 örnek")

    # Sol: tam cümle açıklama
    add_paragraphs(
        s, Inches(0.7), Inches(1.7), Inches(6.5), Inches(4.6),
        [
            "Çalışmada **Kaggle** platformunda yayınlanan **Gameplay Images** "
            "veri seti kullanılmıştır. Veri seti 10 farklı popüler oyundan "
            "**toplam 10.000 ekran görüntüsü** içermektedir; her sınıf için "
            "**1000 örnek** bulunmakta ve sınıf dağılımı **mükemmel şekilde "
            "dengelidir** (sağdaki grafik).",
            "Görseller **640×360 piksel** çözünürlükte PNG formatında "
            "sunulmuştur. Eğitim öncesinde tümü 224×224'e yeniden boyutlandırılarak "
            "ImageNet üzerinde önceden eğitilmiş modellerin beklediği girdi "
            "boyutuyla uyumlu hale getirilmiştir.",
            "Sınıflar: **Among Us, Apex Legends, Fortnite, Forza Horizon, "
            "Free Fire, Genshin Impact, God of War, Minecraft, Roblox, Terraria**.",
        ],
        font_size=12,
    )

    # Sağ: dağılım grafiği
    add_image(s, FIG / "eda_class_distribution.png",
              Inches(7.4), Inches(1.7), width=Inches(5.5))

    add_textbox(
        s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
        "Bu çalışmada veri seti hiçbir sentetik artırma yapılmadan, kaynağında olduğu "
        "gibi 10×1000 dağılımıyla kullanılmıştır.",
        font_size=10, color=SLATE_500,
    )

    apply_chrome(s, 4)


def slide_05_preprocessing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Veri Ön İşleme ve Bölme",
                    "RGBA → RGB · stratified split · augmentation")

    add_paragraphs(
        s, Inches(0.7), Inches(1.65), Inches(12), Inches(2.7),
        [
            "Veri seti üzerindeki ilk işlem, modellerin işleyebileceği tek tip "
            "**RGB formatına dönüştürülmesi** olmuştur. Ham görsellerin yaklaşık "
            "%20'si RGBA (alfa kanallı) formatta olduğundan, eğitim öncesinde "
            "RGB'ye çevrilmiştir. Tüm görseller **224×224 piksele** yeniden "
            "boyutlandırılarak ImageNet pretrained modelleriyle uyumlu hale "
            "getirilmiştir.",
            "Eğitim verisi üzerinde genelleme kabiliyetini artırmak amacıyla "
            "**augmentation** uygulanmıştır: rastgele kırpma (RandomCrop 224), "
            "yatay çevirme, renk titreşimi (ColorJitter — parlaklık, kontrast, "
            "doygunluk) ve ±10° rotation. Validation ve test setlerinde sadece "
            "deterministik resize + center crop kullanılmıştır. ImageNet "
            "ortalama-standart sapma değerleriyle normalizasyon yapılmıştır.",
        ],
        font_size=12,
    )

    # Split kutusu
    add_textbox(s, Inches(0.7), Inches(4.8), Inches(12), Inches(0.4),
                "Veri Bölme — Stratified 70/15/15", font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.6),
        [
            "Veri seti **stratified split** ile bölünmüştür: **%70 eğitim "
            "(7000 görsel), %15 doğrulama (1500 görsel), %15 test (1500 "
            "görsel)**. Stratified yöntemi her sınıftan eşit oranda örnek "
            "alarak train/val/test dağılımını korur — bizim setimiz dengeli "
            "olduğu için sonuç eşit oran ile aynı, ancak yöntem **reproducibility** "
            "ve doğru pratiği sağlar. Bölme **seed=42** ile sabitlenmiş, eğitim "
            "ile test seti hiçbir noktada karışmamıştır.",
        ],
        font_size=12,
    )

    apply_chrome(s, 5)


def slide_06_mlp(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "MLP — Yapay Sinir Ağı Baseline",
                    "Tam bağlı katmanlar · görüntü yapısını yok sayar")

    # Sol: diyagram (büyük)
    draw_mlp_diagram(s, Inches(0.5), Inches(1.65), Inches(7.0), Inches(4.6))

    # Sağ: tam cümle açıklama
    add_paragraphs(
        s, Inches(7.7), Inches(1.65), Inches(5.2), Inches(4.6),
        [
            "MLP baseline'ımız görüntüyü düz bir vektöre çevirerek **üç tam "
            "bağlı katmandan** geçirir. Girdi katmanı 224×224×3 = **150.528** "
            "piksel değeri kabul eder; oradan **256 nöronluk** ilk gizli "
            "katmana, sonra **128 nöronluk** ikinci gizli katmana, son olarak "
            "**10 sınıflık** çıktı katmanına bağlanır.",
            "Her gizli katmandan sonra **ReLU aktivasyonu** ve **Dropout 0.3** "
            "uygulanır. Toplam **38,57 milyon** parametre içerir; bunlar "
            "tamamen tam bağlı katmanlardadır ve **pretrained değildir**.",
            "**Eğitim**: AdamW optimizer (lr=1e-3), CosineAnnealingLR scheduler, "
            "batch size 64, **20 epoch** maksimum, **early stopping patience=5**. "
            "Eğitim **11,8 dakika** sürmüş ve 19. epoch'ta **best val_acc 0,5053** "
            "elde edilmiştir.",
        ],
        font_size=11,
    )

    apply_chrome(s, 6)


def slide_07_cnn_scratch(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "CNN Scratch — Sıfırdan Eğitilen Klasik CNN",
                    "VGG-mini tarzı · 4 evrişim bloğu · 0.42M parametre")

    # Üst: diyagram
    draw_cnn_diagram(s, Inches(0.4), Inches(1.6), Inches(12.5), Inches(2.5))

    # Alt: açıklama
    add_paragraphs(
        s, Inches(0.7), Inches(4.3), Inches(12), Inches(2.5),
        [
            "Sıfırdan eğitilen CNN baseline'ımız **VGG-mini** tarzında **dört "
            "evrişim bloğundan** oluşur. Her blok bir 3×3 convolution, batch "
            "normalization, ReLU aktivasyonu ve 2×2 maxpooling içerir. Filter "
            "sayıları her blokta iki katına çıkar (32 → 64 → 128 → 256), "
            "feature map çözünürlüğü ise yarıya düşer (224 → 112 → 56 → 28 → 14).",
            "Son convolution bloğunun çıktısı **global average pooling** ile "
            "256 boyutlu bir vektöre indirgenir; ardından 128 birimlik bir "
            "tam bağlı katman üzerinden 10 sınıflık tahmin üretilir. Toplam "
            "parametre sayısı yalnızca **0,42 milyondur** — yani MLP'den **92 "
            "kat daha az**, fakat görüntü yapısını sömürdüğü için sonuçları çok "
            "daha iyidir.",
            "**Eğitim**: AdamW (lr=1e-3), AMP mixed precision, **20 epoch** + "
            "**patience=5**. Eğitim **11,1 dakika** sürmüş, 19. epoch'ta "
            "**best val_acc 0,9680** elde edilmiştir — ImageNet pretrained "
            "kullanmadan, sıfırdan.",
        ],
        font_size=11,
    )

    apply_chrome(s, 7)


def slide_08_transfer(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Transfer Learning + 3 Modern Mimari",
                    "ImageNet ön-eğitilmiş ağırlıklar · ResNet, EfficientNet, ViT")

    add_paragraphs(
        s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.0),
        [
            "Transfer learning, ImageNet (1.2 milyon görsel, 1000 sınıf) üzerinde "
            "önceden eğitilmiş ağırlıkların kendi probleme uyarlanmasıdır. Tüm "
            "katmanlar fine-tune edilir; backbone dondurulmaz. Bu strateji, "
            "10K görsellik küçük bir veri setiyle bile **yüksek doğruluk** sağlar.",
        ],
        font_size=12,
    )

    cols = [
        {
            "title": "ResNet50",
            "tag": "2015 · Klasik CNN",
            "params": "23.5M",
            "kernel": "Residual Block",
            "desc": (
                "Skip bağlantılarıyla derin (50+ katman) CNN. "
                "20 epoch + early stop, AMP, lr=1e-4. "
                "10.6 dk'da %99.13 macro-F1."
            ),
        },
        {
            "title": "EfficientNetB0",
            "tag": "2019 · Modern CNN",
            "params": "4.0M",
            "kernel": "MBConv + Compound Scaling",
            "desc": (
                "Mobile inverted bottleneck + 3 boyutu birlikte ölçekler. "
                "20 epoch + early stop, AMP, lr=1e-4. "
                "9.5 dk'da %99.07 macro-F1."
            ),
        },
        {
            "title": "ViT-Base/16",
            "tag": "2020 · Transformer",
            "params": "85.8M",
            "kernel": "Patch Token + Self-Attention",
            "desc": (
                "Görseli 16×16 patch'lere böler, NLP transformer'ını uygular. "
                "20 epoch + early stop, AMP, lr=1e-4, batch 16. "
                "19.5 dk'da %99.07 macro-F1."
            ),
        },
    ]

    col_w = Inches(4.0)
    col_h = Inches(3.5)
    gap = Inches(0.25)
    total = col_w * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    col_y = Inches(2.85)

    for i, c in enumerate(cols):
        x = start_x + (col_w + gap) * i
        add_rounded_rect(s, x, col_y, col_w, col_h, WHITE, GREEN)
        add_rect(s, x, col_y, col_w, Inches(0.55), GREEN_DARK)
        add_textbox(s, x + Inches(0.2), col_y + Inches(0.1),
                    col_w - Inches(0.4), Inches(0.4), c["title"],
                    font_size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, x + Inches(0.2), col_y + Inches(0.7),
                    col_w - Inches(0.4), Inches(0.3), c["tag"],
                    font_size=10, color=SLATE_500)
        add_textbox(s, x + Inches(0.2), col_y + Inches(1.0),
                    col_w - Inches(0.4), Inches(0.4), f"{c['params']} parametre",
                    font_size=12, bold=True, color=GREEN)
        add_textbox(s, x + Inches(0.2), col_y + Inches(1.4),
                    col_w - Inches(0.4), Inches(0.3), "Yapı taşı:",
                    font_size=10, color=SLATE_500)
        add_textbox(s, x + Inches(0.2), col_y + Inches(1.65),
                    col_w - Inches(0.4), Inches(0.35), c["kernel"],
                    font_size=11, bold=True, color=SLATE_900)
        add_textbox(s, x + Inches(0.2), col_y + Inches(2.05),
                    col_w - Inches(0.4), Inches(1.4), c["desc"],
                    font_size=10, color=SLATE_700)

    add_textbox(
        s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
        "Üçü de **timm** kütüphanesi üzerinden ImageNet pretrained ağırlıklarla yüklendi — "
        "tutarlı API ve aynı eğitim transformları sağlar.",
        font_size=10, color=SLATE_500,
    )

    apply_chrome(s, 8)


def slide_09_protocol(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Eğitim Protokolü ve Alternatif Yaklaşımlar",
                    "Adil karşılaştırma için 5 modelde ortak strateji")

    add_paragraphs(
        s, Inches(0.7), Inches(1.65), Inches(6.0), Inches(0.4),
        ["Ortak Hiperparametreler"],
        font_size=13,
    )
    # bullet listesi yerine tam cümle
    add_paragraphs(
        s, Inches(0.7), Inches(2.05), Inches(6.0), Inches(4.5),
        [
            "Tüm modeller **aynı protokolde** eğitildi — adil karşılaştırma için "
            "kritik. Optimizer **AdamW** (weight decay 1e-4), scheduler "
            "**CosineAnnealingLR** ile öğrenme oranı kosinüs eğrisinde 0'a "
            "iniyor. Loss fonksiyonu CrossEntropyLoss.",
            "Maksimum **20 epoch**, **early stopping patience=5** — yani 5 "
            "ardışık epoch iyileşme olmazsa eğitim durur. Baseline modeller "
            "(MLP, CNN scratch) için lr=1e-3, transfer learning modelleri "
            "için daha küçük lr=1e-4 (pretrained ağırlıkları bozmamak için).",
            "**AMP — Automatic Mixed Precision** (FP16 + FP32 karışımı) "
            "transfer ve CNN scratch eğitimlerinde aktif; yaklaşık **%40 "
            "hızlanma** sağladı, doğruluğa zarar vermedi.",
            "Reproducibility: **seed=42**, deterministic dataloader. "
            "Crash recovery için her epoch sonu checkpoint (resume desteği).",
        ],
        font_size=11,
    )

    # Sağ: alternatif yaklaşımlar tablosu
    add_textbox(s, Inches(7.0), Inches(1.65), Inches(5.7), Inches(0.4),
                "Düşünüp Uygulamadığım Alternatifler",
                font_size=13, bold=True, color=GREEN_DARK)

    rows = [
        ("Sıfırdan eğitim", "10K yetmez, %20-30 düşüş bekleniyor"),
        ("Frozen backbone", "Hızlı ama doğruluk %3-5 düşer"),
        ("LoRA / Adapter", "LLM için anlamlı, 86M ViT için marjinal"),
        ("Mixup / CutMix", "%1-2 ek doğruluk ama odağı dağıtır"),
        ("K-Fold CV (5)", "5× hesaplama, tek tahmin için yeterli"),
    ]
    tbl_x = Inches(7.0)
    tbl_y = Inches(2.1)
    col_widths = [Inches(2.3), Inches(3.4)]
    row_h = Inches(0.55)

    headers = ["Yaklaşım", "Neden seçmedim"]
    cur_x = tbl_x
    for hdr, w in zip(headers, col_widths):
        add_rect(s, cur_x, tbl_y, w, Inches(0.4), GREEN_DARK)
        add_textbox(s, cur_x + Inches(0.1), tbl_y, w - Inches(0.2), Inches(0.4),
                    hdr, font_size=11, bold=True, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w
    for i, row in enumerate(rows):
        cur_x = tbl_x
        y = tbl_y + Inches(0.4) + row_h * i
        bg = WHITE if i % 2 == 0 else GREEN_PALE
        for j, (cell, w) in enumerate(zip(row, col_widths)):
            add_rect(s, cur_x, y, w, row_h, bg, SLATE_300)
            add_textbox(s, cur_x + Inches(0.1), y + Inches(0.1),
                        w - Inches(0.2), row_h - Inches(0.2), cell,
                        font_size=10,
                        bold=(j == 0),
                        color=SLATE_900 if j == 0 else SLATE_700,
                        anchor=MSO_ANCHOR.MIDDLE)
            cur_x += w

    apply_chrome(s, 9)


def slide_10_results_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Sonuçlar — 5 Model Test Performansı",
                    "1500 örnekli test seti · birincil metrik Macro-F1")

    headers = ["Model", "Acc", "Macro-F1", "Top-3", "Params", "Boyut", "Inf (ms)", "Eğitim"]
    rows = [
        ("MLP",            "0.5027", "0.4794", "0.7600", "38.57M", "147.1 MB", "9.83",  "11.8 dk"),
        ("CNN Scratch",    "0.9733", "0.9733", "0.9967",  "0.42M",   "1.6 MB", "9.13",  "11.1 dk"),
        ("ResNet50",       "0.9913", "0.9913", "0.9973", "23.53M",  "89.8 MB", "9.29",  "10.6 dk"),
        ("EfficientNetB0", "0.9907", "0.9907", "0.9973",  "4.02M",  "15.3 MB", "9.99",   "9.5 dk"),
        ("ViT-Base/16",    "0.9907", "0.9907", "0.9953", "85.81M", "327.3 MB", "12.01", "19.5 dk"),
    ]

    tbl_x = Inches(0.5)
    tbl_y = Inches(1.7)
    col_w = [Inches(2.5), Inches(1.1), Inches(1.4), Inches(1.0),
             Inches(1.3), Inches(1.4), Inches(1.2), Inches(1.4)]
    row_h = Inches(0.50)
    f1_col_idx = 2  # vurgulanacak

    cur_x = tbl_x
    for i, (hdr, w) in enumerate(zip(headers, col_w)):
        bg = GREEN if i == f1_col_idx else GREEN_DARK
        add_rect(s, cur_x, tbl_y, w, row_h, bg)
        add_textbox(s, cur_x, tbl_y, w, row_h, hdr,
                    font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += w

    total_col_w = Emu(sum(int(w) for w in col_w))
    for ri, row in enumerate(rows):
        cur_x = tbl_x
        y = tbl_y + row_h * (ri + 1)
        # baseline ↔ transfer ayırıcısı
        if ri == 2:
            add_rect(s, tbl_x, _i(y - Inches(0.04)), total_col_w, Inches(0.04), GREEN)
        for ci, (cell, w) in enumerate(zip(row, col_w)):
            bg = WHITE if ri % 2 == 0 else GREEN_PALE
            if ci == f1_col_idx:
                bg = GREEN_PALE  # F1 kolonu hep vurgulu zemin
            add_rect(s, cur_x, y, w, row_h, bg, SLATE_300)
            add_textbox(s, cur_x, y, w, row_h, cell,
                        font_size=11,
                        bold=(ci == 0 or ci == f1_col_idx),
                        color=GREEN_DARK if ci == 0 else (GREEN_DARK if ci == f1_col_idx else SLATE_900),
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cur_x += w

    # Yorum
    add_textbox(s, Inches(0.5), Inches(4.85), Inches(12), Inches(0.35),
                "Anahtar Gözlemler", font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0),
        [
            "**MLP ile transfer arası yaklaşık 0,51 macro-F1 farkı bulunmaktadır** — "
            "görüntü yapısını sömüren mimarilerin değeri net biçimde görülür.",
            "**CNN Scratch yalnızca 0,42M parametre ile 0,9733 macro-F1**'e ulaşmıştır; "
            "ResNet50'nin 0,9913 ile arasında **sadece 0,018 fark** var — yani transfer "
            "learning'in marjinal katkısı bu görevde **küçüktür**, esas sıçrama "
            "MLP→CNN geçişinde yaşanır.",
            "**Transfer learning grubunda fark istatistiksel gürültü içinde** (0,06%); "
            "**EfficientNetB0** 21× daha az parametre ve 6× daha küçük dosyayla "
            "pratik olarak en akıllı tercih.",
        ],
        font_size=11,
    )

    apply_chrome(s, 10)


def slide_11_curves(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Eğitim Eğrileri Karşılaştırması",
                    "5 modelin train/val loss ve accuracy seyri")

    add_image(s, FIG / "comparison_curves_overlay.png",
              Inches(0.4), Inches(1.6), width=Inches(12.5))

    add_rounded_rect(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.7), GREEN_PALE, GREEN)
    add_textbox(
        s, Inches(1.0), Inches(6.42), Inches(11.5), Inches(0.5),
        "MLP düşük accuracy'de plato yapar (%50). CNN Scratch hızlı yakınsar ve %96+ "
        "seviyesine çıkar. Transfer learning modelleri ilk birkaç epoch'ta zaten %95 "
        "üzerine çıkar — ImageNet ağırlıklarının değeri burada görünür.",
        font_size=11, color=SLATE_700,
    )

    apply_chrome(s, 11)


def slide_12_class_f1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Sınıf-Bazlı F1 Analizi",
                    "5 model × 10 sınıf — hangi sınıflar zor?")

    # Sınıf-bazlı F1 bar grup — daha küçük genişlik, ortalanmış
    img_w = Inches(11.0)
    img_h = img_w / 2.5  # ratio 2.5:1
    img_left = (SLIDE_W - img_w) / 2
    add_image(s, FIG / "comparison_class_f1.png", img_left, Inches(1.65), width=img_w)

    # Yorum kutusu — chart'ın altında
    box_top = Inches(1.65) + img_h + Inches(0.15)
    add_rounded_rect(s, Inches(0.5), box_top, Inches(12.3), Inches(1.4),
                     GREEN_PALE, GREEN)
    add_paragraphs(
        s, Inches(0.8), box_top + Inches(0.15), Inches(11.7), Inches(1.2),
        [
            "Her sınıf için 5 modelin F1 skoru yan yana gösterilmektedir "
            "(gri=MLP · turuncu=CNN Scratch · mavi=ResNet50 · yeşil=EfficientNetB0 · "
            "kırmızı=ViT-Base/16).",
            "**MLP**'nin Apex Legends ve Free Fire F1'i 0,2 civarında — model bu iki "
            "FPS'yi birbirine karıştırıyor. **CNN Scratch** ve transfer modelleri tüm "
            "sınıflarda 0,95 üstü; en zor sınıf yine **Apex Legends** (benzer HUD).",
        ],
        font_size=11,
    )

    apply_chrome(s, 12)


def slide_13_confusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Hata Analizi — Confusion Matrices",
                    "5 model için normalize edilmiş confusion matrix grid")

    add_image(s, FIG / "comparison_confusion_grid.png",
              Inches(0.3), Inches(1.65), width=Inches(12.7))

    add_paragraphs(
        s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.6),
        [
            "**Diagonal** = doğru tahmin (recall). **MLP**'nin matrisi "
            "dağınıktır; özellikle Apex Legends, Free Fire ve God of War "
            "satırlarında yanlış tahminler yoğun. **CNN Scratch ve transfer** "
            "modellerinin diagonal'leri belirgin şekilde temizlenir.",
            "Tüm modellerde en zor karışan sınıf ikilisi **Apex Legends ↔ Free Fire** — "
            "her ikisi de battle royale FPS, benzer HUD elementleri (mini harita, sağlık "
            "barı, silah göstergesi) taşıdığı için karıştırılması beklenen bir durum.",
        ],
        font_size=11,
    )

    apply_chrome(s, 13)


def slide_14_gradcam(prs):
    s = prs.slides.add_slide(prs.slide_invitation_layout if False else prs.slide_layouts[6])
    add_slide_title(s, "Grad-CAM ve EigenCAM — Modelin Nereye Baktığı?",
                    "Heatmap görselleştirmesi · neden ViT için EigenCAM?")

    add_paragraphs(
        s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.0),
        [
            "Bir modelin %99 doğruluk vermesi yetmez — **kararına hangi piksellerin "
            "etki ettiği** sorusu da önemlidir. Heatmap görselleştirmesi bu görünmez "
            "mantığı açığa çıkarır; model debug ve güven açısından kritiktir.",
        ],
        font_size=12,
    )

    col_w = Inches(6.0)
    col_y = Inches(2.95)
    col_h = Inches(2.6)

    # GradCAM
    add_rounded_rect(s, Inches(0.7), col_y, col_w, col_h, WHITE, GREEN)
    add_rect(s, Inches(0.7), col_y, col_w, Inches(0.55), GREEN_DARK)
    add_textbox(s, Inches(0.9), col_y + Inches(0.08), col_w - Inches(0.4), Inches(0.4),
                "GradCAM (CNN'ler için)",
                font_size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(
        s, Inches(0.9), col_y + Inches(0.7), col_w - Inches(0.4), col_h - Inches(0.8),
        [
            "Son convolution katmandaki **gradient × activation** çarpımı. "
            "Sınıf-özel: hangi piksel \"Fortnite\" kararına ne kadar katkı verdi.",
            "**CNN Scratch** için son conv bloğu (256 × 14 × 14), **ResNet50** için "
            "layer4'ün son bloğu, **EfficientNetB0** için son MBConv stage'i kullanılır.",
        ],
        font_size=10,
    )

    # EigenCAM
    add_rounded_rect(s, Inches(7.0), col_y, col_w, col_h, WHITE, GREEN)
    add_rect(s, Inches(7.0), col_y, col_w, Inches(0.55), GREEN_DARK)
    add_textbox(s, Inches(7.2), col_y + Inches(0.08), col_w - Inches(0.4), Inches(0.4),
                "EigenCAM (ViT için)",
                font_size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_paragraphs(
        s, Inches(7.2), col_y + Inches(0.7), col_w - Inches(0.4), col_h - Inches(0.8),
        [
            "ViT softmax doygunluğu (conf=1.0) → vanilla GradCAM'in gradient'i tam "
            "**sıfır** olur; heatmap üretilemez.",
            "**EigenCAM** gradient gerektirmez — son block aktivasyonlarının **SVD**'sini "
            "alır. Sınıf-özel değil, ama her görselde **tutarlı** çalışır.",
        ],
        font_size=10,
    )

    add_textbox(
        s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.4),
        "Önemli not: **MLP için heatmap üretilemez** — tam bağlı katmanların "
        "uzamsal feature map'i yoktur.",
        font_size=11, color=SLATE_700,
    )

    add_rounded_rect(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.7), GREEN_PALE, GREEN)
    add_textbox(
        s, Inches(1.0), Inches(6.42), Inches(11.5), Inches(0.5),
        "Pratik anekdot: Demoda Roblox görselinde ViT yanlış cevap verdi; EigenCAM "
        "modelin avatar yerine **gökyüzüne baktığını** gösterdi — hata sebebi anında "
        "anlaşıldı.",
        font_size=10, color=SLATE_700,
    )

    apply_chrome(s, 14)


def slide_15_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Canlı Demo — Lokal Web Uygulaması",
                    "FastAPI + React · 5 model · belirsizlik göstergesi")

    add_paragraphs(
        s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.9),
        [
            "Lokal web uygulaması üç adımdan oluşur: **görsel yükleme** (drag-drop "
            "veya hazır örnekler), **model seçimi** (tek model ya da \"tümünü "
            "karşılaştır\") ve **sonuç görüntüleme** (Top-3 tahminler + Grad-CAM/EigenCAM "
            "heatmap + belirsizlik rozeti).",
        ],
        font_size=12,
    )

    # Belirsizlik göstergesi açıklaması
    add_textbox(s, Inches(0.7), Inches(2.85), Inches(12), Inches(0.4),
                "Belirsizlik Göstergesi (entropy + margin)",
                font_size=14, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.7), Inches(3.25), Inches(12), Inches(2.5),
        [
            "Her tahmin için Shannon **entropy** (sınıf olasılıkları dağılımı) ve "
            "top-1 ile top-2 arasındaki **margin** hesaplanır. Bu iki sinyal **3 "
            "seviyeli rozet** üretir: ",
            "🟢 **Kesin** — yüksek max prob (≥0.85), geniş margin (≥0.50), düşük entropy.",
            "🟡 **Şüpheli** — orta seviye güven; top-1 ve top-2 yarışıyor olabilir.",
            "🔴 **Belirsiz** — düşük max prob (<0.50) veya yüksek entropy.",
            "**OOD sinyali**: comparison modunda 3+ model \"low\" işaretlerse, görsel "
            "muhtemelen 10 sınıf dışı (out-of-distribution).",
        ],
        font_size=11,
    )

    add_rounded_rect(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.95), GREEN_PALE, GREEN)
    add_textbox(
        s, Inches(1.0), Inches(6.13), Inches(11.5), Inches(0.7),
        "▶ Şimdi tarayıcıda **http://localhost:5173** açıyorum: hazır örneklerden bir "
        "konsensüs vakası (Minecraft) ve 3 modelin çeliştiği bir vaka (Fortnite) "
        "göstereceğim, ardından listede olmayan bir oyun yükleyip OOD davranışını "
        "deneyeceğim.",
        font_size=11, bold=True, color=GREEN_DARK,
    )

    apply_chrome(s, 15)


def slide_16_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "Sonuç Değerlendirmesi · Sınırlamalar · Erişim",
                    "Genel değerlendirme · gelecek çalışma · projeye nereden ulaşılır")

    # Sol — sonuç değerlendirmesi
    add_textbox(s, Inches(0.7), Inches(1.7), Inches(6.0), Inches(0.4),
                "Sonuç Değerlendirmesi", font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(0.7), Inches(2.05), Inches(6.0), Inches(2.5),
        [
            "**Mimari fark çoğu zaman maliyet farkıdır.** Bu görevde 0,42M "
            "parametreli CNN Scratch ile 86M parametreli ViT arasında "
            "macro-F1 farkı yalnızca 0,018 — **doğruluk değil, verimlilik "
            "kararı** önemlidir.",
            "**Transfer learning'in en büyük katkısı erken epoch'larda** "
            "görülür; sıfırdan eğitilen CNN ise yeterli veriyle yine "
            "yüksek doğruluğa ulaşır.",
            "**Confidence ≠ Correctness** — model %87 emin olabilir ve "
            "yine de yanılabilir. Belirsizlik göstergesi bu durumu kullanıcıya "
            "anlık olarak gösterir.",
        ],
        font_size=11,
    )

    # Sağ — sınırlamalar
    add_textbox(s, Inches(7.0), Inches(1.7), Inches(5.7), Inches(0.4),
                "Sınırlamalar & Gelecek Çalışma",
                font_size=13, bold=True, color=GREEN_DARK)
    add_paragraphs(
        s, Inches(7.0), Inches(2.05), Inches(5.7), Inches(2.5),
        [
            "**Closed-set softmax** sınır dışı sınıf için zorunlu tahmin "
            "verir; gerçek OOD detection (Mahalanobis, energy-based) eklenmeli.",
            "10K görsel sınırı; daha fazla sınıf ve sahne çeşitliliği "
            "(menü, loading ekranı vb.) eklenmeli.",
            "Confidence calibration (temperature scaling) eksik.",
            "Multi-label genişletme (aynı görselde birden fazla oyun ögesi).",
            "Mobil deploy: EfficientNet → ONNX → mobil cihaz.",
        ],
        font_size=11,
    )

    # Blog erişim şeması
    add_textbox(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
                "Blog Yazısına Erişim", font_size=13, bold=True, color=GREEN_DARK)

    crumbs = [
        "github.com",
        "MAliTopkara",
        "YZM304-Derin-Ogrenme",
        "proje4/",
        "blog.md",
    ]
    crumb_y = Inches(5.15)
    crumb_h = Inches(0.55)
    crumb_w = Inches(2.25)
    arrow_w = Inches(0.18)
    total_w = crumb_w * len(crumbs) + arrow_w * (len(crumbs) - 1)
    start_x = (SLIDE_W - total_w) / 2
    cur_x = start_x

    for i, crumb in enumerate(crumbs):
        col = GREEN if i % 2 == 0 else GREEN_DARK
        add_rounded_rect(s, cur_x, crumb_y, crumb_w, crumb_h, col)
        add_textbox(s, cur_x + Inches(0.1), crumb_y, crumb_w - Inches(0.2), crumb_h,
                    crumb, font_size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += crumb_w
        if i < len(crumbs) - 1:
            arrow = s.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                cur_x, crumb_y + Inches(0.15), arrow_w, crumb_h - Inches(0.3),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GREEN
            arrow.line.fill.background()
            cur_x += arrow_w

    add_textbox(
        s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.4),
        "🔗 https://github.com/MAliTopkara/YZM304-Derin-Ogrenme/tree/main/proje4",
        font_size=11, color=SLATE_500, align=PP_ALIGN.CENTER,
    )

    # References
    add_textbox(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
                "Ana Referanslar", font_size=12, bold=True, color=GREEN_DARK)
    add_textbox(
        s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.5),
        "He et al. 2015 (ResNet) · Tan & Le 2019 (EfficientNet) · "
        "Dosovitskiy et al. 2020 (ViT) · Selvaraju et al. 2017 (Grad-CAM) · "
        "Bany Muhammad & Yeasin 2020 (EigenCAM) · timm (Wightman, HuggingFace)",
        font_size=10, color=SLATE_700,
    )

    apply_chrome(s, 16)


# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01_title,
        slide_02_motivation,
        slide_03_problem,
        slide_04_dataset,
        slide_05_preprocessing,
        slide_06_mlp,
        slide_07_cnn_scratch,
        slide_08_transfer,
        slide_09_protocol,
        slide_10_results_table,
        slide_11_curves,
        slide_12_class_f1,
        slide_13_confusion,
        slide_14_gradcam,
        slide_15_demo,
        slide_16_conclusion,
    ]
    for b in builders:
        b(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"OK: {OUT}  ({len(builders)} slayt)")


if __name__ == "__main__":
    main()
